from pptx import Presentation
from pptx.util import Inches, Pt
import requests
from bs4 import BeautifulSoup
from colors import bcolors
import string

#Global Settings
pTitle = "Careers"
pBy = "Jeff Pilgrim - 3/27/22"
pFileName = "Market Careers.pptx"

#
#links = ["https://www.onetonline.org/link/summary/27-3031.00",
#         "https://www.onetonline.org/link/summary/41-2021.00",
#         "https://www.onetonline.org/link/summary/41-4012.00",
#         "https://www.onetonline.org/link/summary/41-9011.00"]

links = ["http://www.onetonline.org/link/details/23-1011.00#Tasks",
"http://www.onetonline.org/link/details/33-9032.00#Tasks",
"http://www.onetonline.org/link/details/33-3051.01#Tasks",
"http://www.onetonline.org/link/details/33-2011.01#WorkActivities"]

prs = Presentation()

def make_TitleSlide():
    print("Making title slide.")

    #Create slide with properties
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    #Edit Properties
    title.text = pTitle
    subtitle.text = pBy

    print("Completed title slide.")

def addTitle(title_shape, text, category):
    if "," in category:
        category = category.split(",")[0] + " (...)"
    title_shape.text = f"{text} {category}"
    return title_shape

def addBullet(tf, text):
    tf.text = f"{text}"
    return tf

def addBulletP(tf, text):
    tf.text = f"{text}"
    tf.level = 1
    tf.font.size = Pt(12)
    return tf

def refractorInt(st):
    return int(st.replace("$", "").replace(",", ""))

def getRank(list, ann):
    mine = refractorInt(ann)
    list.remove(mine)
    lessthan = 0
    greater = 0
    for num in list:
        if num > mine:
            greater += 1
        if num < mine:
            lessthan += 1
    print("type="+str(mine))
    print("less="+str(lessthan))
    print("greater="+str(greater))
    
    if lessthan > greater:
        if lessthan == 1:
            return "best"
        if lessthan == 3:
            return "best"
        return "best"
    if greater > lessthan:
        return "worst"
    if greater == lessthan:
        return "middle"

def make_MediaSlides():
    print("Making media slides")
    annuals = []
    numbers = []
    for link in links:
        #Common Tasks for blank
        print("Getting attributes from link " + link)
        
        req = requests.get(link)
        html = req.text

        if html.find("Tasks"):
            print("Received HTML")
        else:
            print("Couldn't get HTML. Exiting..")
            return

        #Parse html
        soup = BeautifulSoup(html, 'html.parser')

        #Format the parsed html
        strhtm = soup.prettify()

        #Get title
        rawTitle = str(soup.title)
        split = rawTitle.split(" - ")
        Category = split[1].split("<")[0]
        
        print(f"Got Category - {bcolors.OKGREEN}{Category}{bcolors.ENDC}")

        #Get 4 Tasks
        tasks = []
        classes = [value
                   for element in soup.find_all(class_=True)
                   for value in element["class"]]
        div = soup.find_all('div', {"class" : "moreinfo"})
        for d in div:
            href = d.find("a")
            if "task" in str(href):
                if not "?" in d.text:
                    if not "-" in d.text:
                        if len(tasks) != 4:
                            tasks.insert(len(tasks) + 1, d.text)
                            print(f"Added Task - {bcolors.OKGREEN}{d.text}{bcolors.ENDC}") 
        if len(tasks) != 4:
            classes = [value
                       for element in soup.find_all(class_=True)
                       for value in element["class"]]
            div = soup.find_all('ul', {"class" : "moreinfo"})
            for d in div:
                href = d.find("a")
                if "task" in str(href):
                    if not "?" in d.text:
                        #if not "-" in d.text:
                        if len(tasks) != 4:
                            split = d.text.split(".")
                            for s in split:
                                if len(tasks) != 4:
                                    tasks.insert(len(tasks) + 1, s.strip())
                                    print(f"Added Task - {bcolors.OKGREEN}{s.strip()}{bcolors.ENDC}")
        if len(tasks) != 4:
            classes = [value
                       for element in soup.find_all(class_=True)
                       for value in element["class"]]
            div = soup.find_all('div', {"class" : "d-flex flex-nowrap pb-1"})
            for d in div:
                href = d.find("a")
                newtext = d.find('div', {"class" : "order-2 flex-grow-1"})
                if "task" in str(href):
                    if not "?" in newtext.text:
                        #if not "-" in d.text:
                        if len(tasks) != 4:
                            split = newtext.text.split(".")
                            for s in split:
                                if len(tasks) != 4 and s.strip():
                                    tasks.insert(len(tasks) + 1, s.strip())
                                    print(f"Added Task - {bcolors.OKGREEN}{s.strip()}{bcolors.ENDC}")
        
        if len(tasks) != 4:
            print(f"Only contains {bcolors.FAIL}{len(tasks)}{bcolors.ENDC} out of 4 tasks. Exiting..")
            return
            
        #Get 4 Work Activities
        activities = []
        classes = [value
                   for element in soup.find_all(class_=True)
                   for value in element["class"]]
        div = soup.find_all('div', {"class" : "moreinfo"})
        for d in div:
            href = d.find("a")
            if "workactivities" in str(href):
                if not "?" in d.text:
                    if "-" in d.text:
                        if len(activities) != 4:
                            activities.insert(len(activities) + 1, d.text)
                            print(f"Added Activity - {bcolors.OKGREEN}{d.text}{bcolors.ENDC}")
        if len(activities) != 4:
            classes = [value
                       for element in soup.find_all(class_=True)
                       for value in element["class"]]
            div = soup.find_all('ul', {"class" : "moreinfo"})
            for d in div:
                href = d.find("a")
                if "workactivities" in str(href):
                    if not "?" in d.text:
                        if "-" in d.text:
                            if len(activities) != 4:
                                split = d.text.split(".")
                                for s in split:
                                    if len(activities) != 4:
                                        activities.insert(len(activities) + 1, s.strip())
                                        print(f"Added Activity - {bcolors.OKGREEN}{s.strip()}{bcolors.ENDC}")
        if len(activities) != 4:
            classes = [value
                       for element in soup.find_all(class_=True)
                       for value in element["class"]]
            div = soup.find_all('div', {"class" : "d-flex flex-nowrap pb-1"})
            for d in div:
                href = d.find("a")
                newtext = d.find('div', {"class" : "order-2 flex-grow-1"})
                if "workactivities" in str(href):
                    if not "?" in newtext.text:
                        if "-" in newtext.text:
                            if len(activities) != 4:
                                split = newtext.text.split(".")
                                for s in split:
                                    if len(activities) != 4 and s.strip():
                                        activities.insert(len(activities) + 1, s.strip())
                                        print(f"Added Activity - {bcolors.OKGREEN}{s.strip()}{bcolors.ENDC}")
        if len(activities) != 4:
            print(f"Only contains {bcolors.FAIL}{len(activities)}{bcolors.ENDC} out of 4 activities. Exiting..")
            return
            
        #Get 4 Work Contexts
        contexts = []
        classes = [value
                   for element in soup.find_all(class_=True)
                   for value in element["class"]]
        div = soup.find_all('div', {"class" : "moreinfo"})
        for d in div:
            href = d.find("a")
            if "workcontext" in str(href):
                if "?" in d.text:
                    if "-" in d.text:
                        if len(contexts) != 4:
                            contexts.insert(len(contexts) + 1, d.text)
                            print(f"Added Context - {bcolors.OKGREEN}{d.text}{bcolors.ENDC}")
        if len(contexts) != 4:
            classes = [value
                       for element in soup.find_all(class_=True)
                       for value in element["class"]]
            div = soup.find_all('ul', {"class" : "moreinfo"})
            for d in div:
                href = d.find("a")
                if "workcontext" in str(href):
                    if not "?" in d.text:
                        if "-" in d.text:
                            if len(contexts) != 4:
                                split = d.text.split(".")
                                for s in split:
                                    if len(contexts) != 4:
                                        #l = s.split("—")[0].strip().rstrip().lstrip().split(" ")
                                        #print(string.ascii_letters)
                                        #filtered = [s for s in l if all(c in string.ascii_letters for c in s)]
                                        #s = ""
                                        #for xx in filtered:
                                        #    s+= xx + " "
                                        contexts.insert(len(contexts) + 1, s.replace("”", "").replace("“", "").strip())
                                        print(f"Added Context - {bcolors.OKGREEN}{s.replace('”', '').replace('“', '').strip()}{bcolors.ENDC}")
        if len(contexts) != 4: #and not "%" in newtext.text and not "responded" in newtext.text
            classes = [value
                       for element in soup.find_all(class_=True)
                       for value in element["class"]]
            div = soup.find_all('div', {"class" : "d-flex flex-nowrap pb-1"})
            for d in div:
                href = d.find("a")
                newtext = d.find('div', {"class" : "order-2 flex-grow-1"})
                if "workcontext" in str(href):
                    if "?" in newtext.text:
                        if "-" in newtext.text:
                            if len(contexts) != 4:
                                split = newtext.text.split(".")
                                for s in split:
                                    if len(contexts) != 4 and s.strip():
                                        if "%" in s.strip():
                                            contexts.insert(len(contexts) + 1, s.strip().split("?")[0] + "?")
                                            print(f"Added Context - {bcolors.OKGREEN}{s.strip().split('?')[0] + '?'}{bcolors.ENDC}")
                                        else:
                                            contexts.insert(len(contexts) + 1, s.strip())
                                            print(f"Added Context - {bcolors.OKGREEN}{s.strip()}{bcolors.ENDC}")
        if len(contexts) != 4:
            print(f"Only contains {bcolors.FAIL}{len(contexts)}{bcolors.ENDC} out of 4 contexts. Exiting..")
            return
            
        #Make tasks slide
        
        for num in range(0,3):
            print(f"Current index {bcolors.WARNING}{num}{bcolors.ENDC}")
            bullet_slide_layout = prs.slide_layouts[1]

            slide = prs.slides.add_slide(bullet_slide_layout)
            shapes = slide.shapes

            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            tf = body_shape.text_frame
            
            if num == 0:
                addTitle(title_shape, "Common Tasks for", Category)
                addBullet(tf, "Tasks")
                addBulletP(tf.add_paragraph(), "1. " + tasks[1] + "\n2. " + tasks[2] + "\n3. " + tasks[3] + "\n4. " + tasks[0])
                
            if num == 1:
                addTitle(title_shape, "Common Activities for", Category)
                addBullet(tf, "Activities")
                addBulletP(tf.add_paragraph(), "1. " + activities[1] + "\n2. " + activities[2] + "\n3. " + activities[3] + "\n4. " + activities[0])
                
            if num == 2:
                addTitle(title_shape, "Common Contexts for", Category)
                addBullet(tf, "Contexts")
                addBulletP(tf.add_paragraph(), "1. " + contexts[1] + "\n2. " + contexts[2] + "\n3. " + contexts[3] + "\n4. " + contexts[0])
        #Get annual salary
              
        annual = []
        classes = [value
                   for element in soup.find_all(class_=True)
                   for value in element["class"]]
        td = soup.find_all('td', {"class" : "report2"})
        for d in td:
            if "$" in d.text:
                if "annual" in d.text:
                    a = d.text.split("hourly, ")[1].split(" annual")[0]
                    annuals.insert(len(annuals) + 1, Category + "|" + a)
                    numbers.insert(len(numbers) + 1, refractorInt(a))
                    print(f"Added Annual - {bcolors.OKGREEN}{a}{bcolors.ENDC}")
              
        if len(annuals) != 4:
            classes = [value
                       for element in soup.find_all(class_=True)
                       for value in element["class"]]
            td = soup.find_all('dd', {"class" : "col-sm-9 col-form-label pt-xso-0"})
            for d in td:
                if "$" in d.text:
                    if "annual" in d.text:
                        a = d.text.split("hourly, ")[1].split(" annual")[0]
                        annuals.insert(len(annuals) + 1, Category + "|" + a)
                        numbers.insert(len(numbers) + 1, refractorInt(a))
                        print(f"Added Annual - {bcolors.OKGREEN}{a}{bcolors.ENDC}")
            
    
    print(f"{bcolors.WARNING}Annuals = ")
    for s in annuals:
        print(s)
    print(bcolors.ENDC)
    
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    
    shapes.title.text = "Career Comparison"
    cols = 3
    rows = 5
    left = Inches(1.25)
    top = Inches(2.25)
    width = Inches(6.0)
    height = Inches(0.8)
    
    table = shapes.add_table(rows, cols, left, top, width, height).table

    # set column widths
    table.columns[0].width = Inches(4.0)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(1.5)

    # write column headings
    table.cell(0, 0).text = 'Career'
    table.cell(0, 1).text = 'Income'
    table.cell(0, 2).text = 'Rank'

    # write body cells
    table.cell(1, 0).text = annuals[0].split("|")[0]
    table.cell(1, 1).text = annuals[0].split("|")[1]
    table.cell(1, 2).text = getRank(numbers, annuals[0].split("|")[1])

    
    table.cell(2, 0).text = annuals[1].split("|")[0]
    table.cell(2, 1).text = annuals[1].split("|")[1]
    table.cell(2, 2).text = getRank(numbers, annuals[1].split("|")[1])
    
    table.cell(3, 0).text = annuals[2].split("|")[0]
    table.cell(3, 1).text = annuals[2].split("|")[1]
    table.cell(3, 2).text = getRank(numbers, annuals[2].split("|")[1])
    
    table.cell(4, 0).text = annuals[3].split("|")[0]
    table.cell(4, 1).text = annuals[3].split("|")[1]
    table.cell(4, 2).text = getRank(numbers, annuals[3].split("|")[1])
    
    #Add Works Cited
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    tf = body_shape.text_frame
    addTitle(title_shape, "Works Cited", "")
    addBullet(tf, "Sources")
    #Save presentation
    for link in links:
        addBulletP(tf.add_paragraph(), link)
        
    save_Slides()

def save_Slides():
    print("Saving Presentation")
    
    prs.save(pFileName)

    print("Saved Presentation as " + pFileName)

make_TitleSlide()

make_MediaSlides()




